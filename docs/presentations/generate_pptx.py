#!/usr/bin/env python3
"""
PowerPoint Generator for Contract-First Integration Presentation
Creates a modern, professional presentation from slides-content.md

Author: Wallace Espindola
Email: wallace.espindola@gmail.com
"""

import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pathlib import Path


class ModernPresentationGenerator:
    """Generates a modern PowerPoint presentation with custom theming"""

    # Modern color palette (based on technical theme)
    COLORS = {
        'primary': RGBColor(0, 102, 204),      # Blue
        'secondary': RGBColor(255, 87, 34),    # Orange accent
        'dark': RGBColor(33, 33, 33),          # Almost black
        'light': RGBColor(250, 250, 250),      # Off-white
        'code_bg': RGBColor(40, 44, 52),       # Dark code background
        'success': RGBColor(76, 175, 80),      # Green
        'warning': RGBColor(255, 193, 7),      # Yellow
        'error': RGBColor(244, 67, 54),        # Red
    }

    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(16)
        self.prs.slide_height = Inches(9)

    def add_title_slide(self, title, subtitle, author_info):
        """Create modern title slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank layout

        # Add gradient background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.COLORS['primary']

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(2.5), Inches(14), Inches(2)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(60)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)
        title_para.alignment = PP_ALIGN.CENTER

        # Subtitle
        subtitle_box = slide.shapes.add_textbox(
            Inches(1), Inches(4.5), Inches(14), Inches(1)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.font.size = Pt(32)
        subtitle_para.font.color.rgb = RGBColor(255, 255, 255)
        subtitle_para.alignment = PP_ALIGN.CENTER

        # Author info
        author_box = slide.shapes.add_textbox(
            Inches(1), Inches(6.5), Inches(14), Inches(1.5)
        )
        author_frame = author_box.text_frame
        author_frame.text = author_info
        author_para = author_frame.paragraphs[0]
        author_para.font.size = Pt(18)
        author_para.font.color.rgb = RGBColor(255, 255, 255)
        author_para.alignment = PP_ALIGN.CENTER

    def add_content_slide(self, title, content):
        """Create content slide with modern styling"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.COLORS['light']

        # Title bar with accent color
        title_shape = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0), Inches(0),
            Inches(16), Inches(1.2)
        )
        title_fill = title_shape.fill
        title_fill.solid()
        title_fill.fore_color.rgb = self.COLORS['primary']

        # Title text
        title_frame = title_shape.text_frame
        title_frame.text = title
        title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(36)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)
        title_para.alignment = PP_ALIGN.LEFT

        # Content area
        content_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.5), Inches(15), Inches(7)
        )
        content_frame = content_box.text_frame
        content_frame.word_wrap = True

        # Parse and add content
        self._add_formatted_content(content_frame, content)

        # Footer with page indicator
        self._add_footer(slide, len(self.prs.slides))

    def add_code_slide(self, title, code, language=""):
        """Create slide with syntax-highlighted code"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.COLORS['light']

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(15), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = self.COLORS['dark']

        # Code box with dark background
        code_shape = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0.5), Inches(1.3),
            Inches(15), Inches(7.2)
        )
        code_fill = code_shape.fill
        code_fill.solid()
        code_fill.fore_color.rgb = self.COLORS['code_bg']

        # Code text
        code_frame = code_shape.text_frame
        code_frame.text = code
        code_para = code_frame.paragraphs[0]
        code_para.font.name = 'Consolas'
        code_para.font.size = Pt(16)
        code_para.font.color.rgb = RGBColor(171, 178, 191)  # Light gray for code

        # Language indicator
        if language:
            lang_box = slide.shapes.add_textbox(
                Inches(13.5), Inches(1.0), Inches(2), Inches(0.4)
            )
            lang_frame = lang_box.text_frame
            lang_frame.text = language.upper()
            lang_para = lang_frame.paragraphs[0]
            lang_para.font.size = Pt(12)
            lang_para.font.bold = True
            lang_para.font.color.rgb = self.COLORS['secondary']
            lang_para.alignment = PP_ALIGN.RIGHT

        self._add_footer(slide, len(self.prs.slides))

    def add_section_slide(self, title, subtitle=""):
        """Create section divider slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Background with gradient effect (simulated with two rectangles)
        bg_shape1 = slide.shapes.add_shape(
            1, Inches(0), Inches(0), Inches(16), Inches(9)
        )
        bg_fill1 = bg_shape1.fill
        bg_fill1.solid()
        bg_fill1.fore_color.rgb = self.COLORS['primary']

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(3), Inches(14), Inches(2)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(54)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)
        title_para.alignment = PP_ALIGN.CENTER

        if subtitle:
            subtitle_box = slide.shapes.add_textbox(
                Inches(1), Inches(5.5), Inches(14), Inches(1)
            )
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.text = subtitle
            subtitle_para = subtitle_frame.paragraphs[0]
            subtitle_para.font.size = Pt(28)
            subtitle_para.font.color.rgb = RGBColor(255, 255, 255)
            subtitle_para.alignment = PP_ALIGN.CENTER

    def _add_formatted_content(self, text_frame, content):
        """Add formatted content with bullets, emphasis, etc."""
        lines = content.strip().split('\n')

        for i, line in enumerate(lines):
            if i > 0:
                p = text_frame.add_paragraph()
            else:
                p = text_frame.paragraphs[0]

            # Check for different formatting
            if line.startswith('- ') or line.startswith('• '):
                # Bullet point
                p.text = line[2:].strip()
                p.level = 0
                p.font.size = Pt(22)
                p.font.color.rgb = self.COLORS['dark']
                p.space_before = Pt(6)

            elif line.startswith('  - '):
                # Sub-bullet
                p.text = line[4:].strip()
                p.level = 1
                p.font.size = Pt(20)
                p.font.color.rgb = self.COLORS['dark']

            elif line.startswith('**') and line.endswith('**'):
                # Bold emphasis
                p.text = line.strip('*')
                p.font.size = Pt(24)
                p.font.bold = True
                p.font.color.rgb = self.COLORS['primary']
                p.space_before = Pt(12)

            elif line.startswith('✅') or line.startswith('❌') or line.startswith('🔹'):
                # Icon lines
                p.text = line
                p.font.size = Pt(22)
                p.font.color.rgb = self.COLORS['dark']
                p.space_before = Pt(8)

            elif line.strip():
                # Regular text
                p.text = line
                p.font.size = Pt(22)
                p.font.color.rgb = self.COLORS['dark']
                p.space_before = Pt(6)

    def _add_footer(self, slide, page_num):
        """Add footer with page number and branding"""
        footer_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(8.3), Inches(15), Inches(0.4)
        )
        footer_frame = footer_box.text_frame
        footer_frame.text = f"Contract-First Integration  |  Wallace Espindola  |  Slide {page_num}"
        footer_para = footer_frame.paragraphs[0]
        footer_para.font.size = Pt(12)
        footer_para.font.color.rgb = RGBColor(150, 150, 150)
        footer_para.alignment = PP_ALIGN.CENTER

    def save(self, filename):
        """Save the presentation"""
        self.prs.save(filename)
        print(f"✓ Presentation saved: {filename}")


def parse_markdown_slides(md_content):
    """Parse markdown content into slides"""
    slides = []
    current_slide = {"title": "", "content": "", "type": "content"}

    lines = md_content.split('\n')
    i = 0

    while i < len(lines):
        line = lines[i].strip()

        # Slide separator
        if line == '---':
            if current_slide["title"] or current_slide["content"]:
                slides.append(current_slide)
            current_slide = {"title": "", "content": "", "type": "content"}
            i += 1
            continue

        # Title (H1 or H2)
        if line.startswith('# ') and not current_slide["title"]:
            current_slide["title"] = line[2:].strip()
            current_slide["type"] = "title" if i == 0 else "content"
            i += 1
            continue

        if line.startswith('## '):
            if current_slide["title"] or current_slide["content"]:
                slides.append(current_slide)
                current_slide = {"title": "", "content": "", "type": "content"}
            current_slide["title"] = line[3:].strip()
            i += 1
            continue

        # Code block
        if line.startswith('```'):
            language = line[3:].strip()
            code_lines = []
            i += 1
            while i < len(lines) and not lines[i].strip().startswith('```'):
                code_lines.append(lines[i])
                i += 1
            current_slide["content"] = '\n'.join(code_lines)
            current_slide["type"] = "code"
            current_slide["language"] = language
            i += 1
            continue

        # Regular content
        if line:
            current_slide["content"] += line + '\n'

        i += 1

    # Add last slide
    if current_slide["title"] or current_slide["content"]:
        slides.append(current_slide)

    return slides


def main():
    """Generate PowerPoint presentation from markdown"""
    print("=" * 60)
    print("Contract-First Integration - PowerPoint Generator")
    print("=" * 60)
    print()

    # Read markdown content
    md_file = Path(__file__).parent / "slides-content.md"
    if not md_file.exists():
        print(f"✗ Error: {md_file} not found")
        return 1

    print(f"Reading slides content from: {md_file}")
    with open(md_file, 'r', encoding='utf-8') as f:
        md_content = f.read()

    # Parse slides
    print("Parsing slides...")
    slides_data = parse_markdown_slides(md_content)
    print(f"✓ Found {len(slides_data)} slides")
    print()

    # Generate presentation
    print("Generating PowerPoint presentation...")
    gen = ModernPresentationGenerator()

    for i, slide_data in enumerate(slides_data, 1):
        print(f"  Creating slide {i}/{len(slides_data)}: {slide_data['title'][:50]}...")

        if slide_data["type"] == "title" and i == 1:
            # Title slide
            gen.add_title_slide(
                title=slide_data["title"],
                subtitle=slide_data["content"].split('\n')[0] if slide_data["content"] else "",
                author_info=slide_data["content"]
            )
        elif slide_data["type"] == "code":
            # Code slide
            gen.add_code_slide(
                title=slide_data["title"],
                code=slide_data["content"],
                language=slide_data.get("language", "")
            )
        elif "Key Takeaways" in slide_data["title"] or "Questions" in slide_data["title"]:
            # Section divider
            gen.add_section_slide(
                title=slide_data["title"],
                subtitle=""
            )
        else:
            # Regular content slide
            gen.add_content_slide(
                title=slide_data["title"],
                content=slide_data["content"]
            )

    # Save presentation
    output_file = Path(__file__).parent / "contract-first-integration.pptx"
    gen.save(str(output_file))

    print()
    print("=" * 60)
    print("✓ Presentation generated successfully!")
    print(f"  Location: {output_file}")
    print(f"  Slides: {len(slides_data)}")
    print()
    print("Next steps:")
    print("  1. Open in PowerPoint, Google Slides, or Keynote")
    print("  2. Customize design and animations")
    print("  3. Add speaker notes if needed")
    print("  4. Present!")
    print("=" * 60)

    return 0


if __name__ == "__main__":
    import sys
    sys.exit(main())
